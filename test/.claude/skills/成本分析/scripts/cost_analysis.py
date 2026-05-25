#!/usr/bin/env python3
"""
新华医院成本分析模块 PPT 数据更新工具。

用法:
    python cost_analysis.py <excel_path> <pptx_path> [--output OUTPUT]
    python cost_analysis.py 4月分析底稿.xlsx 新华医院2026年4月运营分析0509.pptx

从 Excel 附表和数据表中提取关键成本/收入数据，除以 10000 后更新到 PPT 成本分析模块。
"""

import json
import os
import sys
import re
import argparse
from datetime import datetime
from collections import defaultdict

import openpyxl
from pptx import Presentation


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def detect_data_month(excel_path):
    """Extract month number from filename like '4月分析底稿.xlsx' → 4."""
    basename = os.path.basename(excel_path)
    m = re.search(r'(\d{1,2})\s*月', basename)
    if m:
        return int(m.group(1))
    return 4  # default


# ---------------------------------------------------------------------------
# ExcelDataLoader
# ---------------------------------------------------------------------------

class ExcelDataLoader:
    """加载 Excel 所有 sheet，建立 {source_alias: {normalized_label: {col_idx: value}}} 索引."""

    def __init__(self, excel_path, sources_config):
        self.excel_path = excel_path
        self.sources = sources_config
        self.cache = {}       # source_alias -> {label -> {col: value}}
        self.needs_wan = {}   # source_alias -> bool
        self.label_col = {}   # source_alias -> int
        self._load_all()

    @staticmethod
    def _norm(s):
        """标准化文本用于匹配."""
        if s is None:
            return ""
        s = str(s).strip()
        s = re.sub(r'[\s　\n\r\t]+', '', s)
        return s

    def _match_sheet(self, wb, target_name):
        """通过精确匹配 → 标准化匹配 → 包含匹配找 sheet."""
        # Exact match
        for ws in wb.worksheets:
            if ws.title == target_name:
                return ws
        # Normalized match
        tn = self._norm(target_name)
        for ws in wb.worksheets:
            if self._norm(ws.title) == tn:
                return ws
        # Substring either way
        for ws in wb.worksheets:
            wsn = self._norm(ws.title)
            if tn and (tn in wsn or wsn in tn):
                return ws
        return None

    def _load_all(self):
        wb = openpyxl.load_workbook(self.excel_path, data_only=True)
        available = [ws.title for ws in wb.worksheets]
        print(f"  Excel sheets ({len(available)}): {available}")

        for alias, cfg in self.sources.items():
            sheet_name = cfg["sheet"]
            needs = cfg.get("needs_wan", False)
            self.needs_wan[alias] = needs
            label_col = cfg.get("label_col", 0)
            self.label_col[alias] = label_col

            ws = self._match_sheet(wb, sheet_name)
            if ws is None:
                print(f"  [WARN] Sheet not found: '{sheet_name}'")
                self.cache[alias] = {}
                continue

            sheet_data = {}
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
                label = self._norm(row[label_col].value)
                if label:
                    row_dict = {}
                    for c, cell in enumerate(row):
                        row_dict[c] = cell.value
                    sheet_data[label] = row_dict
            self.cache[alias] = sheet_data
            print(f"  [LOAD] {alias} ← {ws.title} ({len(sheet_data)} rows)")
        wb.close()

    def _fuzzy_get(self, sheet_data, target_label):
        """模糊匹配：exact > contains > startswith."""
        if target_label in sheet_data:
            return sheet_data[target_label]
        for label in sheet_data:
            if target_label in label:
                return sheet_data[label]
        for label in sheet_data:
            if label in target_label:
                return sheet_data[label]
        return None

    def get_value(self, source_alias, row_label, col_idx):
        """获取指定 source 中指定行标签、列的值，自动处理万元转换."""
        if source_alias not in self.cache:
            return None
        sheet_data = self.cache[source_alias]
        row = self._fuzzy_get(sheet_data, self._norm(row_label))
        if row is None:
            return None
        actual_col = col_idx + self.label_col.get(source_alias, 0)
        val = row.get(actual_col)
        if val is None:
            return None
        if not isinstance(val, (int, float)):
            return val
        if self.needs_wan.get(source_alias, False):
            val = val / 10000.0
        return val

    def get_raw_value(self, source_alias, row_label, col_idx):
        """Like get_value but without wan conversion. Used for sum computations."""
        if source_alias not in self.cache:
            return None
        sheet_data = self.cache[source_alias]
        row = self._fuzzy_get(sheet_data, self._norm(row_label))
        if row is None:
            return None
        actual_col = col_idx + self.label_col.get(source_alias, 0)
        val = row.get(actual_col)
        if val is None:
            return None
        if not isinstance(val, (int, float)):
            return val
        return val


# ---------------------------------------------------------------------------
# MaterialDetailLoader
# ---------------------------------------------------------------------------

class MaterialDetailLoader:
    """加载卫生材料明细 Excel，提供按条件聚合查询能力.

    明细文件有两个 sheet，列布局不同：
    - 2026 sheet (20列): col 7=院区, col 16=是否收费, col 12=入账科目, col 3=部门名称, col 19=金额
    - 2025 sheet (19列): 无院区列, col 15=是否收费, col 11=入账科目, col 3=部门名称, col 18=金额
    """

    def __init__(self, detail_path):
        self.detail_path = detail_path
        self.data = {}  # year -> list of dicts
        self._load()

    def _load(self):
        wb = openpyxl.load_workbook(self.detail_path, data_only=True)
        for ws in wb.worksheets:
            if ws.max_column < 2:
                continue
            ncols = ws.max_column
            if ncols == 20:
                year = 2026
                col_map = {"period": 1, "campus": 7, "chargeable": 16, "subject": 12, "dept": 3, "amount": 19}
            elif ncols == 19:
                year = 2025
                col_map = {"period": 1, "chargeable": 15, "subject": 11, "dept": 3, "amount": 18}
            else:
                continue

            rows = []
            for r in range(2, ws.max_row + 1):
                period_raw = str(ws.cell(r, col_map["period"] + 1).value or "").strip()
                chargeable_raw = str(ws.cell(r, col_map["chargeable"] + 1).value or "").strip()
                subject = str(ws.cell(r, col_map["subject"] + 1).value or "").strip()
                dept = str(ws.cell(r, col_map["dept"] + 1).value or "").strip()
                amount = ws.cell(r, col_map["amount"] + 1).value
                if not isinstance(amount, (int, float)):
                    continue
                row_data = {
                    "period": period_raw,
                    "chargeable": chargeable_raw == "是",
                    "subject": subject,
                    "dept": dept,
                    "amount": amount,
                }
                if "campus" in col_map:
                    row_data["campus"] = str(ws.cell(r, col_map["campus"] + 1).value or "").strip()
                rows.append(row_data)

            self.data[year] = rows
            campus_info = "with campus" if "campus" in col_map else "no campus"
            print(f"  [LOAD] 卫生材料明细 {year} ← {ws.title} ({len(rows)} rows, {campus_info})")
        wb.close()

    def aggregate(self, year, chargeable=None, subject_contains=None,
                  dept_contains=None, campus=None, period=None):
        """按条件过滤并返回金额总和（元）.

        Args:
            year: 2026 or 2025
            chargeable: True=是, False=否, None=不限
            subject_contains: 入账科目包含此字符串
            dept_contains: 部门名称包含此字符串
            campus: 院区（仅2026数据有效）
            period: 期间过滤，如 "2026-04"（None=不限）

        Returns:
            金额总和（元），若无匹配返回 0.
        """
        if year not in self.data:
            return 0
        total = 0.0
        for row in self.data[year]:
            if period and row["period"] != period:
                continue
            if chargeable is not None and row["chargeable"] != chargeable:
                continue
            if subject_contains and subject_contains not in row["subject"]:
                continue
            if dept_contains and dept_contains not in row["dept"]:
                continue
            if campus and row.get("campus") and campus not in row["campus"]:
                continue
            total += row["amount"]
        return total


# ---------------------------------------------------------------------------
# ComputedValueResolver
# ---------------------------------------------------------------------------

class ComputedValueResolver:
    """Resolves computed column and text variable specifications."""

    def __init__(self, loader, data_month=None, detail_loader=None):
        self.loader = loader
        self.data_month = data_month or 4
        self.detail_loader = detail_loader

    def resolve(self, compute_spec, resolved_cols=None):
        ctype = compute_spec["type"]
        if ctype == "value":
            return compute_spec["value"]
        elif ctype == "subtract":
            a = resolved_cols[compute_spec["a_col"]]
            b = resolved_cols[compute_spec["b_col"]]
            return a - b if a is not None and b is not None else None
        elif ctype == "add":
            a = resolved_cols[compute_spec["a_col"]]
            b = resolved_cols[compute_spec["b_col"]]
            return a + b if a is not None and b is not None else None
        elif ctype == "divide":
            a = resolved_cols[compute_spec["a_col"]]
            b = resolved_cols[compute_spec["b_col"]]
            if a is not None and b is not None and b != 0:
                return a / b
            return None
        elif ctype == "multiply":
            a = resolved_cols[compute_spec["a_col"]]
            b = resolved_cols[compute_spec["b_col"]]
            return a * b if a is not None and b is not None else None
        elif ctype == "sum":
            total = 0
            for part in compute_spec["parts"]:
                v = self.loader.get_raw_value(part["source"], part.get("xl_label", ""), part["xl_col"])
                if v is not None and isinstance(v, (int, float)):
                    total += v
                else:
                    return None
            if compute_spec.get("needs_wan", False):
                total /= 10000.0
            return total
        elif ctype == "time_node_diff":
            exec_rate = resolved_cols[compute_spec["exec_rate_col"]]
            if exec_rate is None:
                return None
            month = compute_spec.get("month", self.data_month)
            return exec_rate - (month / 12.0)
        elif ctype == "pct_diff":
            a = resolved_cols[compute_spec["a_col"]]
            b = resolved_cols[compute_spec["b_col"]]
            if a is not None and b is not None and b != 0:
                return (a - b) / abs(b)
            return None
        elif ctype == "pct_change":
            change = resolved_cols[compute_spec["change_col"]]
            base = resolved_cols[compute_spec["base_col"]]
            if change is not None and base is not None:
                denominator = base - change
                if denominator != 0:
                    return change / abs(denominator)
            return None
        elif ctype == "divide_excel":
            a = self._resolve_source(compute_spec["a"])
            b = self._resolve_source(compute_spec["b"])
            if a is not None and b is not None and b != 0:
                return a / b
            return None
        elif ctype == "time_node_diff_excel":
            exec_rate = None
            if "exec_rate" in compute_spec:
                exec_rate = self._resolve_source(compute_spec["exec_rate"])
            elif "a" in compute_spec and "b" in compute_spec:
                a = self._resolve_source(compute_spec["a"])
                b = self._resolve_source(compute_spec["b"])
                if a is not None and b is not None and b != 0:
                    exec_rate = a / b
            if exec_rate is None:
                return None
            month = compute_spec.get("month", self.data_month)
            return exec_rate - (month / 12.0)
        elif ctype == "subtract_excel":
            a = self._resolve_source(compute_spec["a"])
            b = self._resolve_source(compute_spec["b"])
            return a - b if a is not None and b is not None else None
        elif ctype == "sum_excel":
            total = 0
            for part in compute_spec["parts"]:
                # Use get_raw_value to avoid double-dividing by 10000
                v = self.loader.get_raw_value(part["source"], part.get("xl_label", ""), part["xl_col"])
                if v is not None and isinstance(v, (int, float)):
                    total += v
                else:
                    return None
            if compute_spec.get("needs_wan", False):
                total /= 10000.0
            return total
        elif ctype == "detail_sum":
            if self.detail_loader is None:
                return None
            year = compute_spec["year"]
            month = compute_spec.get("month", self.data_month)
            period = f"{year}-{month:02d}"
            raw = self.detail_loader.aggregate(
                year=year,
                chargeable=compute_spec.get("chargeable"),
                subject_contains=compute_spec.get("subject_contains"),
                dept_contains=compute_spec.get("dept_contains"),
                campus=compute_spec.get("campus"),
                period=period,
            )
            if compute_spec.get("needs_wan", False):
                raw = raw / 10000.0
            return raw
        elif ctype == "add_refs":
            a = self._resolve_source(compute_spec["a"])
            b = self._resolve_source(compute_spec["b"])
            return a + b if a is not None and b is not None else None
        elif ctype == "subtract_refs":
            a = self._resolve_source(compute_spec["a"])
            b = self._resolve_source(compute_spec["b"])
            return a - b if a is not None and b is not None else None
        return None

    def _resolve_source(self, source_spec):
        """Resolve a value from an Excel source spec or compute spec."""
        if isinstance(source_spec, dict):
            if "compute" in source_spec:
                return self.resolve(source_spec["compute"])
            return self.loader.get_value(
                source_spec.get("source", ""),
                source_spec.get("xl_label", ""),
                source_spec.get("xl_col", 0)
            )
        return source_spec


# ---------------------------------------------------------------------------
# TextUpdater
# ---------------------------------------------------------------------------

class TextUpdater:
    """Updates text paragraphs in non-table shapes (text boxes)."""

    @staticmethod
    def norm(s):
        if s is None:
            return ""
        return re.sub(r'[\s　\n\r]+', '', str(s).strip())

    @classmethod
    def find_shape_by_name(cls, slide, shape_name):
        target = cls.norm(shape_name)
        for shape in slide.shapes:
            if cls.norm(shape.name) == target:
                return shape
        return None

    @staticmethod
    def format_var(value, format_type):
        if value is None:
            return "N/A"
        if not isinstance(value, (int, float)):
            return str(value)
        if format_type == "wan_int":
            return f"{round(value / 10000.0):,}"
        elif format_type == "wan_1dp":
            return f"{value / 10000.0:,.1f}"
        elif format_type == "pct_1dp":
            return f"{value * 100:.1f}%"
        elif format_type == "pct_2dp":
            return f"{value * 100:.2f}%"
        elif format_type == "pp_1dp":
            return f"{abs(value) * 100:.1f}"
        elif format_type == "pct_abs_1dp":
            return f"{abs(value) * 100:.1f}%"
        elif format_type == "int":
            return f"{round(value):,}"
        else:
            if isinstance(value, float) and abs(value) < 10 and value != int(value):
                return f"{value:.2f}"
            return f"{round(value):,}"

    @classmethod
    def update_text(cls, slide, text_config, loader, resolver):
        shape_name = text_config["shape_name"]
        shape = cls.find_shape_by_name(slide, shape_name)
        if shape is None:
            return False

        var_values = {}
        for var_name, var_spec in text_config.get("vars", {}).items():
            if var_spec.get("type") == "value":
                var_values[var_name] = var_spec["value"]
            elif "compute" in var_spec:
                var_values[var_name] = resolver.resolve(var_spec["compute"])
            else:
                val = loader.get_value(
                    var_spec.get("source", ""),
                    var_spec.get("xl_label", ""),
                    var_spec.get("xl_col", 0)
                )
                var_values[var_name] = val

        format_overrides = text_config.get("format_overrides", {})
        formatted_values = {}
        for var_name, val in var_values.items():
            fmt = format_overrides.get(var_name)
            if fmt:
                formatted_values[var_name] = cls.format_var(val, fmt)
            else:
                if isinstance(val, (int, float)):
                    formatted_values[var_name] = f"{round(val):,}"
                else:
                    formatted_values[var_name] = str(val) if val is not None else "N/A"

        template = text_config["template"]
        new_text = template.format(**formatted_values)

        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = new_text
            # Copy font properties if original had them
            return True

        return False


# ---------------------------------------------------------------------------
# PPTTableFinder
# ---------------------------------------------------------------------------

class PPTTableFinder:

    @staticmethod
    def norm(s):
        if s is None:
            return ""
        return re.sub(r'[\s　\n\r]+', '', str(s).strip())

    @classmethod
    def find_table_by_header(cls, slide, headers, header_row=0):
        """在 slide 中查找包含指定表头的表格."""
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            if header_row >= len(table.rows):
                continue
            row_cells = [cls.norm(table.rows[header_row].cells[c].text) for c in range(len(table.columns))]
            # Check if all headers appear as substrings in the row
            match_count = 0
            for h in headers:
                h_norm = cls.norm(h)
                for cell_text in row_cells:
                    if h_norm and h_norm in cell_text:
                        match_count += 1
                        break
            if match_count >= len(headers) * 0.6:  # 60% threshold
                return table
        return None

    @classmethod
    def find_row_by_label(cls, table, label, col_idx=0):
        """在表格第 col_idx 列中查找包含 label 的行号."""
        target = cls.norm(label)
        for r_idx, row in enumerate(table.rows):
            cell_text = cls.norm(row.cells[col_idx].text)
            if target and target in cell_text:
                return r_idx
        # Try looser match
        for r_idx, row in enumerate(table.rows):
            cell_text = cls.norm(row.cells[col_idx].text)
            if target and cell_text and cell_text in target:
                return r_idx
        return None

    @classmethod
    def get_row_labels(cls, table, col_idx=0):
        """获取表格指定列的所有行标签."""
        return [cls.norm(table.rows[r].cells[col_idx].text) for r in range(len(table.rows))]


# ---------------------------------------------------------------------------
# ValueFormatter
# ---------------------------------------------------------------------------

class ValueFormatter:
    """格式化更新到 PPT 单元格的值."""

    @staticmethod
    def is_percentage_text(val):
        """判断原始文本是否包含百分号."""
        if isinstance(val, str):
            return '%' in val
        return False

    @staticmethod
    def is_decimal(val):
        """判断值是否为小数（非金额）."""
        if isinstance(val, float) and abs(val) < 100 and val != int(val):
            return True
        return False

    @classmethod
    def format_value(cls, raw_value, ppt_original_text, format_override=None):
        """根据 PPT 原始格式和新值，决定输出文本.

        Args:
            raw_value: 从 Excel 获取的新值（数值或字符串）.
            ppt_original_text: PPT 单元格原始文本，用于判断格式.
            format_override: 可选的格式字符串，如 "2dp", "pct", "int".

        Returns:
            格式化后的字符串.
        """
        if raw_value is None:
            return ppt_original_text  # 保持原样

        # 如果新值是字符串（如百分比），直接返回
        if isinstance(raw_value, str):
            return raw_value

        if not isinstance(raw_value, (int, float)):
            return str(raw_value)

        # Explicit format override takes priority
        if format_override == "2dp":
            return f"{raw_value:,.2f}"
        elif format_override == "int":
            return f"{round(raw_value):,}"

        # 判断原始 PPT 单元格的格式
        ppt_text = str(ppt_original_text).strip() if ppt_original_text else ""

        # 包含 % → 按百分比格式化（Excel 百分比数据均为比值形式：0.3459→34.6%, 1→100%, 1.1875→118.8%）
        if '%' in ppt_text:
            return f"{raw_value * 100:.1f}%"

        # 负值 → 带符号
        val = raw_value

        # 根据原始文本判断是否需要小数
        if re.search(r'\.\d', ppt_text):
            # 原始有小数 → 保留小数
            return f"{val:,.2f}"
        else:
            # 原始为整数 → 四舍五入
            rounded = round(val)
            return f"{rounded:,}"


# ---------------------------------------------------------------------------
# CostAnalysisUpdater
# ---------------------------------------------------------------------------

class CostAnalysisUpdater:
    """主控制器：根据 mapping 更新 PPT."""

    def __init__(self, excel_path, pptx_path, mapping_path):
        self.excel_path = excel_path
        self.pptx_path = pptx_path
        self.mapping_path = mapping_path

        with open(mapping_path, 'r', encoding='utf-8') as f:
            self.mapping = json.load(f)

        self.loader = ExcelDataLoader(excel_path, self.mapping.get("excel_sources", {}))
        self.prs = Presentation(pptx_path)
        self.finder = PPTTableFinder()
        self.formatter = ValueFormatter()

        self.data_month = detect_data_month(excel_path)

        # Auto-detect 卫生材料明细 file in same directory as main Excel
        self.detail_loader = None
        excel_dir = os.path.dirname(os.path.abspath(excel_path))
        for fname in os.listdir(excel_dir):
            if '卫生材料明细' in fname and fname.endswith('.xlsx') and not fname.startswith('~$'):
                detail_path = os.path.join(excel_dir, fname)
                print(f"  [DETAIL] Found: {fname}")
                self.detail_loader = MaterialDetailLoader(detail_path)
                break
        if self.detail_loader is None:
            print("  [DETAIL] No detail file found, detail_sum compute type will return None")

        self.resolver = ComputedValueResolver(self.loader, self.data_month, self.detail_loader)
        self.text_updater = TextUpdater()

        self.stats = {"updated_cells": 0, "skipped_rows": 0, "skipped_cells": 0, "slides_processed": 0}

    def run(self, output_path):
        slides_config = self.mapping.get("slides", {})
        total_slides = len(slides_config)

        for slide_key, slide_cfg in slides_config.items():
            slide_num = int(slide_key)
            if slide_num < 1 or slide_num > len(self.prs.slides):
                print(f"  [SKIP] Slide {slide_num}: out of range")
                continue

            slide = self.prs.slides[slide_num - 1]  # 1-indexed → 0-indexed
            print(f"\n--- Slide {slide_num}: {slide_cfg.get('desc', '')} ---")
            self._process_slide(slide, slide_cfg)

        self.prs.save(output_path)
        print(f"\n{'='*50}")
        print(f"Summary: Updated {self.stats['updated_cells']} cells across "
              f"{self.stats['slides_processed']} slides, "
              f"skipped {self.stats['skipped_rows']} rows, "
              f"{self.stats['skipped_cells']} cells.")
        print(f"Output: {output_path}")

    def _process_slide(self, slide, slide_cfg):
        headers = slide_cfg.get("headers", [])
        header_row = slide_cfg.get("header_row", 0)

        # Phase 1: Table updates
        table = self.finder.find_table_by_header(slide, headers, header_row)
        if table is None:
            print(f"  [SKIP] Table not found with headers: {headers}")
        else:
            self.stats["slides_processed"] += 1
            rows_cfg = slide_cfg.get("rows", [])
            for row_cfg in rows_cfg:
                self._process_table_row(table, row_cfg, slide_cfg)

        # Phase 2: Text paragraph updates (independent of table existence)
        texts_cfg = slide_cfg.get("texts", [])
        for text_cfg in texts_cfg:
            success = self.text_updater.update_text(slide, text_cfg, self.loader, self.resolver)
            if success:
                print(f"  [UPD] Text shape '{text_cfg['shape_name']}' updated")
            else:
                print(f"  [SKIP] Text shape '{text_cfg['shape_name']}' not found")

    def _process_table_row(self, table, row_cfg, slide_cfg):
        label = row_cfg["label"]
        default_source = row_cfg.get("source")
        default_xl_label = row_cfg.get("xl_label", label)

        row_idx = self.finder.find_row_by_label(table, label, col_idx=0)
        if row_idx is None:
            print(f"  [SKIP Row] '{label}' not found in table")
            self.stats["skipped_rows"] += 1
            return

        col_entries = row_cfg.get("cols", [])
        resolved = {}   # {ppt_col: value}
        col_fmts = {}   # {ppt_col: format_override}

        # Pass 1: Resolve all simple (non-computed) column values
        for entry in col_entries:
            if isinstance(entry, list):
                ppt_col, xl_col = entry[0], entry[1]
                val = self.loader.get_value(default_source, default_xl_label, xl_col)
                resolved[ppt_col] = val
            elif isinstance(entry, dict):
                if "compute" in entry:
                    continue  # Defer to pass 2
                else:
                    ppt_col = entry["ppt"]
                    src = entry.get("source", default_source)
                    xlbl = entry.get("xl_label", default_xl_label)
                    xl_col = entry.get("xl", 0)
                    val = self.loader.get_value(src, xlbl, xl_col)
                    resolved[ppt_col] = val
                    if "format" in entry:
                        col_fmts[ppt_col] = entry["format"]

        # Pass 2: Evaluate computed columns (in order, so dependencies must be earlier)
        for entry in col_entries:
            if isinstance(entry, dict) and "compute" in entry:
                ppt_col = entry["ppt"]
                val = self.resolver.resolve(entry["compute"], resolved_cols=resolved)
                resolved[ppt_col] = val
                if "format" in entry:
                    col_fmts[ppt_col] = entry["format"]

        # Apply all resolved values to table cells
        for ppt_col, new_val in resolved.items():
            if new_val is None:
                print(f"  [SKIP Cell] Row='{label}', PPT col={ppt_col}: no value resolved")
                self.stats["skipped_cells"] += 1
                continue

            cell = table.rows[row_idx].cells[ppt_col]
            original_text = cell.text
            fmt = col_fmts.get(ppt_col)
            formatted = self.formatter.format_value(new_val, original_text, format_override=fmt)

            if formatted != original_text:
                self._update_cell_text(cell, formatted)
                self.stats["updated_cells"] += 1
                print(f"  [UPD] [{row_idx},{ppt_col}] '{label}': \"{original_text[:20]}\" → \"{formatted}\"")

    @staticmethod
    def _update_cell_text(cell, new_text):
        """Write new text to a table cell, replacing existing content."""
        if cell.text_frame.paragraphs:
            for para in cell.text_frame.paragraphs:
                if para.runs:
                    para.runs[0].text = new_text
                else:
                    para.add_run().text = new_text
                break
        else:
            cell.text_frame.paragraphs[0].add_run().text = new_text


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="新华医院成本分析 PPT 数据更新工具")
    parser.add_argument("excel", help="Excel 底稿文件路径")
    parser.add_argument("pptx", help="PPT 模板文件路径")
    parser.add_argument("--output", "-o", default=None, help="输出路径（默认在原 PPT 名加日期后缀）")
    parser.add_argument("--mapping", "-m", default=None, help="mapping.json 路径")
    args = parser.parse_args()

    # 验证输入文件
    for path, desc in [(args.excel, "Excel"), (args.pptx, "PPTX")]:
        if not os.path.exists(path):
            print(f"[ERROR] {desc} file not found: {path}")
            sys.exit(1)

    # 默认 mapping 路径
    if args.mapping is None:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        args.mapping = os.path.join(script_dir, "..", "references", "mapping.json")

    if not os.path.exists(args.mapping):
        print(f"[ERROR] Mapping file not found: {args.mapping}")
        sys.exit(1)

    # 默认输出路径
    if args.output is None:
        base, ext = os.path.splitext(args.pptx)
        date_str = datetime.now().strftime("%m%d")
        args.output = f"{base}_updated_{date_str}{ext}"

    print(f"Excel:  {args.excel}")
    print(f"PPTX:   {args.pptx}")
    print(f"Output: {args.output}")
    print(f"Mapping: {args.mapping}")
    print()

    updater = CostAnalysisUpdater(args.excel, args.pptx, args.mapping)
    updater.run(args.output)


if __name__ == "__main__":
    main()
